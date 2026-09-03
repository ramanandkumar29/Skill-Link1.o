import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Premium Color Palette
    BG_DARK = RGBColor(10, 16, 30)         # Slate 950 #0A101E
    CARD_BG = RGBColor(22, 33, 54)         # Slate 900 Card #162136
    CARD_BORDER = RGBColor(45, 62, 90)     # Slate 700
    
    PRIMARY_BLUE = RGBColor(37, 99, 235)   # Blue 600 #2563EB
    CYAN_ACCENT = RGBColor(56, 189, 248)   # Sky 400 #38BDF8
    SOS_RED = RGBColor(239, 68, 68)        # Red 500 #EF4444
    AMBER_GOLD = RGBColor(245, 158, 11)    # Amber 500 #F59E0B
    EMERALD_GREEN = RGBColor(16, 185, 129) # Emerald 500 #10B981
    PURPLE_ACCENT = RGBColor(168, 85, 247) # Purple 500 #A855F7
    PINK_ACCENT = RGBColor(236, 72, 153)   # Pink 500 #EC4899

    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_MUTED = RGBColor(148, 163, 184)   # Slate 400

    def set_slide_background(slide, color):
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = color
        bg_shape.line.fill.background()
        return bg_shape

    def add_header(slide, title_text, category_text="SKILL-LINK PLATFORM OVERVIEW", character_badge=None):
        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.5), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = "⚡ " + category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = CYAN_ACCENT

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(9.5), Inches(0.65))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(20)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE

        # Character Badge on right if specified
        if character_badge:
            badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.2), Inches(0.55), Inches(2.33), Inches(0.45))
            badge.fill.solid()
            badge.fill.fore_color.rgb = RGBColor(30, 41, 59)
            badge.line.color.rgb = CYAN_ACCENT
            badge.line.width = Pt(1)
            tf_b = badge.text_frame
            p_b = tf_b.paragraphs[0]
            p_b.text = character_badge
            p_b.font.size = Pt(9)
            p_b.font.bold = True
            p_b.font.color.rgb = TEXT_WHITE
            p_b.alignment = PP_ALIGN.CENTER

        # Accent Line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.733), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = CARD_BORDER
        line.line.fill.background()

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
        return card

    # =========================================================================
    # SLIDE 1: Title Slide & Mascot Introduction
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1, BG_DARK)

    # Hero Card Container
    add_card(slide1, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9), CARD_BG, CARD_BORDER)

    # Badge Pill
    badge1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.2), Inches(3.8), Inches(0.36))
    badge1.fill.solid()
    badge1.fill.fore_color.rgb = RGBColor(30, 58, 138)
    badge1.line.color.rgb = CYAN_ACCENT
    tf_b1 = badge1.text_frame
    p_b1 = tf_b1.paragraphs[0]
    p_b1.text = "⚡ NEXT-GEN HYPERLOCAL PLATFORM"
    p_b1.font.size = Pt(9)
    p_b1.font.bold = True
    p_b1.font.color.rgb = CYAN_ACCENT

    # Main Title
    t_box1 = slide1.shapes.add_textbox(Inches(1.2), Inches(1.65), Inches(6.8), Inches(1.1))
    tf1 = t_box1.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "Skill-Link"
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE

    # Subtitle
    sub_box1 = slide1.shapes.add_textbox(Inches(1.2), Inches(2.75), Inches(6.8), Inches(0.7))
    tf_sub1 = sub_box1.text_frame
    tf_sub1.word_wrap = True
    p_sub1 = tf_sub1.paragraphs[0]
    p_sub1.text = "AI-Powered Skilled Worker & 15-Min Roadside Emergency Ecosystem"
    p_sub1.font.size = Pt(14)
    p_sub1.font.bold = True
    p_sub1.font.color.rgb = CYAN_ACCENT

    # Desc
    desc_box1 = slide1.shapes.add_textbox(Inches(1.2), Inches(3.5), Inches(6.8), Inches(1.4))
    tf_d1 = desc_box1.text_frame
    tf_d1.word_wrap = True
    p_d1 = tf_d1.paragraphs[0]
    p_d1.text = "Bridging verified local service professionals, 15-minute emergency roadside mechanics, and everyday citizens through autonomous AI reasoning, multilingual voice UX, and transparent pricing."
    p_d1.font.size = Pt(11)
    p_d1.font.color.rgb = TEXT_MUTED

    # Persona Showcase Cards on the right
    # Persona 1: Priya
    p_card1 = add_card(slide1, Inches(8.3), Inches(1.2), Inches(3.8), Inches(1.5), RGBColor(30, 41, 59), PINK_ACCENT)
    tb_p1 = slide1.shapes.add_textbox(Inches(8.5), Inches(1.3), Inches(3.4), Inches(1.3))
    tf_pc1 = tb_p1.text_frame
    tf_pc1.word_wrap = True
    p_h1 = tf_pc1.paragraphs[0]
    p_h1.text = "👩‍💼 Priya (Customer / Commuter)"
    p_h1.font.size = Pt(11)
    p_h1.font.bold = True
    p_h1.font.color.rgb = PINK_ACCENT
    p_b1_desc = tf_pc1.add_paragraph()
    p_b1_desc.text = "“Stranded with a flat tyre at 9 PM on the highway. Needs 1-tap rapid help without extortion.”"
    p_b1_desc.font.size = Pt(9.5)
    p_b1_desc.font.color.rgb = TEXT_MUTED

    # Persona 2: Lexi AI
    p_card2 = add_card(slide1, Inches(8.3), Inches(2.9), Inches(3.8), Inches(1.5), RGBColor(30, 41, 59), CYAN_ACCENT)
    tb_p2 = slide1.shapes.add_textbox(Inches(8.5), Inches(3.0), Inches(3.4), Inches(1.3))
    tf_pc2 = tb_p2.text_frame
    tf_pc2.word_wrap = True
    p_h2 = tf_pc2.paragraphs[0]
    p_h2.text = "🤖 Lexi AI & Sahayak Voice"
    p_h2.font.size = Pt(11)
    p_h2.font.bold = True
    p_h2.font.color.rgb = CYAN_ACCENT
    p_b2_desc = tf_pc2.add_paragraph()
    p_b2_desc.text = "“Parses Hindi voice queries, searches RAG vector database, and executes autonomous tool dispatches.”"
    p_b2_desc.font.size = Pt(9.5)
    p_b2_desc.font.color.rgb = TEXT_MUTED

    # Persona 3: Vikram
    p_card3 = add_card(slide1, Inches(8.3), Inches(4.6), Inches(3.8), Inches(1.5), RGBColor(30, 41, 59), AMBER_GOLD)
    tb_p3 = slide1.shapes.add_textbox(Inches(8.5), Inches(4.7), Inches(3.4), Inches(1.3))
    tf_pc3 = tb_p3.text_frame
    tf_pc3.word_wrap = True
    p_h3 = tf_pc3.paragraphs[0]
    p_h3.text = "👷‍♂️ Vikram (Verified Pro Mechanic)"
    p_h3.font.size = Pt(11)
    p_h3.font.bold = True
    p_h3.font.color.rgb = AMBER_GOLD
    p_b3_desc = tf_pc3.add_paragraph()
    p_b3_desc.text = "“Accepts SOS ticket in 6s, arrives on motorcycle in 12 mins with tools, and earns fair rates with 0% cut.”"
    p_b3_desc.font.size = Pt(9.5)
    p_b3_desc.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 2: Problem Statement Storyboard
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2, BG_DARK)
    add_header(slide2, "The Real-World Dilemma: Roadside Emergencies & Service Chaos", "PROBLEM STATEMENT", "👩‍💼 Priya's Crisis")

    probs = [
        ("🚨 Highway Stranding", "Vehicle breakdowns, flat tyres, or dead batteries leave commuters stranded on dark highways with zero verified rapid help.", SOS_RED, "Avg Wait: 2-4 Hours"),
        ("💸 Unregulated Pricing", "Finding reliable mechanics relies on luck. Unverified contractors demand arbitrary surge pricing and extortionate cash fees.", AMBER_GOLD, "Markup: Up to 300%"),
        ("🗣️ Language Barrier", "Over 65% of Bharat users struggle with typing English search terms, creating high friction without conversational voice UX.", CYAN_ACCENT, "Exclusion: 65% Users"),
        ("🛡️ Trust Deficit", "Lack of background verification, transparent rate cards, authentic peer ratings, or guarantee safeguards creates anxiety.", PURPLE_ACCENT, "Safety Recourse: 0%")
    ]

    for i, (title, desc, color, stat) in enumerate(probs):
        left = Inches(0.8 + i * 3.0)
        card = add_card(slide2, left, Inches(1.8), Inches(2.75), Inches(4.8), CARD_BG, color)
        
        # Pill on top
        pill = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.2), Inches(2.0), Inches(2.35), Inches(0.32))
        pill.fill.solid()
        pill.fill.fore_color.rgb = RGBColor(30, 41, 59)
        pill.line.color.rgb = color
        p_tf = pill.text_frame
        p_p = p_tf.paragraphs[0]
        p_p.text = stat
        p_p.font.size = Pt(8.5)
        p_p.font.bold = True
        p_p.font.color.rgb = color
        p_p.alignment = PP_ALIGN.CENTER

        # Content Box
        tbox = slide2.shapes.add_textbox(left + Inches(0.2), Inches(2.5), Inches(2.35), Inches(3.8))
        tf = tbox.text_frame
        tf.word_wrap = True
        
        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(13)
        pt.font.bold = True
        pt.font.color.rgb = TEXT_WHITE
        
        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(10)
        pd.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 3: The Solution Ecosystem
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3, BG_DARK)
    add_header(slide3, "Skill-Link: The Intelligent Hyperlocal Ecosystem", "VALUE PROPOSITION", "🤖 Lexi AI Bridge")

    # Left Hub Card
    add_card(slide3, Inches(0.8), Inches(1.8), Inches(4.5), Inches(4.8), CARD_BG, PRIMARY_BLUE)
    tbox_l = slide3.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(3.9), Inches(4.3))
    tf_l = tbox_l.text_frame
    tf_l.word_wrap = True
    
    pl1 = tf_l.paragraphs[0]
    pl1.text = "🎯 Unified Marketplace Solution"
    pl1.font.size = Pt(16)
    pl1.font.bold = True
    pl1.font.color.rgb = CYAN_ACCENT

    bullets_l = [
        "Dual-Mode Operational Engine: Seamless toggle between scheduled home repairs and emergency SOS.",
        "Voice-First Vernacular UX: Bilingual Hindi and English natural speech interactions powered by Web Speech API.",
        "Lexi Autonomous RAG Agent: Instant worker matchmaking and structured tool calling in under 80ms.",
        "Zero-Middlemen Direct Booking: Connects verified workers directly with customers at 0-5% commission."
    ]
    for b in bullets_l:
        p = tf_l.add_paragraph()
        p.text = "➔ " + b
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_MUTED

    # Right 3 Pillar Cards
    pillars = [
        ("🚨 15-Minute Priority Roadside SOS", "Live browser GPS coordinates lock + auto-dispatching mobile mechanics for tyre punctures, dead batteries, and engine stalls with a 15-second ticket auto-lock.", SOS_RED),
        ("🛡️ Verified Trust & Fixed Rate Cards", "Background-checked professionals, authentic peer ratings, transparent upfront rate cards, and an automatic 7-day post-service guarantee.", EMERALD_GREEN),
        ("🛠️ Self-Service QuickFix DIY & Brand Directory", "Interactive diagnostic trees for guided troubleshooting + verified official toll-free anti-fraud directory for 50+ auto & home appliance brands.", AMBER_GOLD)
    ]
    for i, (title, desc, color) in enumerate(pillars):
        top = Inches(1.8 + i * 1.65)
        add_card(slide3, Inches(5.6), top, Inches(6.9), Inches(1.45), CARD_BG, color)
        tb = slide3.shapes.add_textbox(Inches(5.8), top + Inches(0.12), Inches(6.5), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(12)
        pt.font.bold = True
        pt.font.color.rgb = color
        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(9.5)
        pd.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 4: 15-Minute Roadside SOS Feature
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4, BG_DARK)
    add_header(slide4, "Feature Deep-Dive: 15-Minute Priority On-Road Emergency SOS", "CORE CAPABILITIES", "👷‍♂️ Vikram Dispatched")

    # Left: Emergency Pipeline
    add_card(slide4, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), CARD_BG, SOS_RED)
    tb4_l = slide4.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.3))
    tf4_l = tb4_l.text_frame
    tf4_l.word_wrap = True
    p4_1 = tf4_l.paragraphs[0]
    p4_1.text = "🚨 Highway Emergency Flow"
    p4_1.font.size = Pt(15)
    p4_1.font.bold = True
    p4_1.font.color.rgb = SOS_RED

    steps_sos = [
        ("1. Geolocation Lock", "HTML5 High-Accuracy GPS pins user coordinates on highway bypass in 1.2 seconds."),
        ("2. Breakdown Triage", "User picks issue: Tyre Puncture, Battery Jumpstart, Fuel Delivery, Towing, or Engine Breakdown."),
        ("3. 15s Auto-Lock Dispatch", "System pings top 3 nearest mobile mechanics; first to accept locks the dispatch instantly."),
        ("4. Live En-Route Tracking", "Customer views live ETA countdown and one-tap direct phone coordinate call.")
    ]
    for s_title, s_desc in steps_sos:
        p = tf4_l.add_paragraph()
        p.text = "⚡ " + s_title + ": " + s_desc
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_WHITE

    # Right: Home Services Comparison
    add_card(slide4, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), CARD_BG, PRIMARY_BLUE)
    tb4_r = slide4.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.3))
    tf4_r = tb4_r.text_frame
    tf4_r.word_wrap = True
    p4_2 = tf4_r.paragraphs[0]
    p4_2.text = "🏡 Hyperlocal Home Services"
    p4_2.font.size = Pt(15)
    p4_2.font.bold = True
    p4_2.font.color.rgb = CYAN_ACCENT

    steps_home = [
        ("Multi-Category Marketplace", "Electricians, Plumbers, Carpenters, Painters, AC Technicians, and Appliance Repair."),
        ("Verified Worker Profiles", "Ratings, verified skill badges, total completed jobs, and verified peer reviews."),
        ("Slot & Date Scheduling", "Flexible morning/afternoon appointment slots with upfront transparent pricing."),
        ("Direct WhatsApp Coordination", "Zero middlemen interference with 7-day post-service warranty.")
    ]
    for s_title, s_desc in steps_home:
        p = tf4_r.add_paragraph()
        p.text = "➔ " + s_title + ": " + s_desc
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_WHITE

    # =========================================================================
    # SLIDE 5: AI & Multilingual Voice Architecture
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5, BG_DARK)
    add_header(slide5, "Conversational AI: Sahayak Voice & Lexi Autonomous RAG", "AI & VOICE INNOVATION", "🎙️ Hindi + English")

    # Left: Sahayak Voice Engine
    add_card(slide5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), CARD_BG, EMERALD_GREEN)
    tb5_l = slide5.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.3))
    tf5_l = tb5_l.text_frame
    tf5_l.word_wrap = True
    p5_l = tf5_l.paragraphs[0]
    p5_l.text = "🎙️ Sahayak Multilingual Voice Engine"
    p5_l.font.size = Pt(15)
    p5_l.font.bold = True
    p5_l.font.color.rgb = EMERALD_GREEN

    voice_pts = [
        ("Zero-Typing UX", "Empowers elderly and non-tech-savvy users to speak natural colloquial queries."),
        ("Hinglish & Vernacular Parsing", "Seamlessly extracts intent from Hindi phrases like 'पंखा ठीक करने वाला चाहिए'."),
        ("Interactive Speech Synthesis", "Speaks confirmation responses aloud in native natural accents."),
        ("Live Waveform Visualizer", "Dynamic audio pulse indicators give active listening feedback.")
    ]
    for vt, vd in voice_pts:
        p = tf5_l.add_paragraph()
        p.text = "➔ " + vt + ": " + vd
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_WHITE

    # Right: Lexi RAG Agent
    add_card(slide5, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), CARD_BG, PURPLE_ACCENT)
    tb5_r = slide5.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.3))
    tf5_r = tb5_r.text_frame
    tf5_r.word_wrap = True
    p5_r = tf5_r.paragraphs[0]
    p5_r.text = "🤖 Lexi Autonomous RAG Agent"
    p5_r.font.size = Pt(15)
    p5_r.font.bold = True
    p5_r.font.color.rgb = PURPLE_ACCENT

    rag_pts = [
        ("OpenRouter LLM Reasoning", "High-capacity reasoning engine for intent classification and entity slot filling."),
        ("RAG Vector Indexing", "Indexes rate cards, neighborhood coverage, and safety warranty policies."),
        ("Autonomous Function Calling", "Directly executes searchWorkers() and bookEmergencyMechanic() backend tools."),
        ("Embedded UI Message Cards", "Chat returns actionable booking widgets and map buttons directly inside conversation.")
    ]
    for rt, rd in rag_pts:
        p = tf5_r.add_paragraph()
        p.text = "➔ " + rt + ": " + rd
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_WHITE

    # =========================================================================
    # SLIDE 6: QuickFix DIY & Brand Helplines
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide6, BG_DARK)
    add_header(slide6, "Value-Add Utilities: QuickFix Interactive DIY & Brand Directory", "ECOSYSTEM UTILITIES", "🛡️ Anti-Fraud Directory")

    # Left: QuickFix DIY
    add_card(slide6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), CARD_BG, AMBER_GOLD)
    tb6_l = slide6.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.3))
    tf6_l = tb6_l.text_frame
    tf6_l.word_wrap = True
    p6_l = tf6_l.paragraphs[0]
    p6_l.text = "🛠️ QuickFix Step-by-Step DIY Assistant"
    p6_l.font.size = Pt(15)
    p6_l.font.bold = True
    p6_l.font.color.rgb = AMBER_GOLD

    qf_pts = [
        ("Guided Decision Trees", "Self-solve common minor issues (tripping MCBs, leaking taps, AC filter wash)."),
        ("Prominent Safety Warnings", "Mandatory safety check alerts before handling live electrical lines or gas valves."),
        ("Tool & Difficulty Checklist", "Estimated time, tool requirements, and complexity rating before starting."),
        ("1-Tap Pro Escalation", "Instant handoff button to assign a nearby technician if the DIY repair is too difficult.")
    ]
    for qt, qd in qf_pts:
        p = tf6_l.add_paragraph()
        p.text = "➔ " + qt + ": " + qd
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_WHITE

    # Right: Verified Directory
    add_card(slide6, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), CARD_BG, CYAN_ACCENT)
    tb6_r = slide6.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.3))
    tf6_r = tb6_r.text_frame
    tf6_r.word_wrap = True
    p6_r = tf6_r.paragraphs[0]
    p6_r.text = "📞 50+ Verified Brand Toll-Free Directory"
    p6_r.font.size = Pt(15)
    p6_r.font.bold = True
    p6_r.font.color.rgb = CYAN_ACCENT

    dir_pts = [
        ("Anti-Fraud Protection", "Curated official customer care directory protecting users from fake Google ad scams."),
        ("Automotive OEM Support", "Maruti Suzuki SOS, Hyundai Roadside, Tata Motors 24x7, Mahindra Emergency."),
        ("Appliance Brand Support", "Voltas AC Care, Havells, Samsung, LG, Daikin, Godrej, Whirlpool, Crompton."),
        ("1-Touch Direct Dialing", "Instant click-to-call integration for rapid warranty service claims.")
    ]
    for dt, dd in dir_pts:
        p = tf6_r.add_paragraph()
        p.text = "➔ " + dt + ": " + dd
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_WHITE

    # =========================================================================
    # SLIDE 7: Full-Stack Architecture
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide7, BG_DARK)
    add_header(slide7, "Full-Stack System Architecture & Technology Stack", "TECHNICAL ARCHITECTURE", "⚡ Modern Stack")

    tiers = [
        ("🎨 Presentation Layer", [
            "Next.js 14 App Router: Server-side rendering and client routing",
            "React 18 & TypeScript: Type-safe modular component state",
            "TailwindCSS: Dark glassmorphic responsive design system",
            "Web Speech API: Low-latency client STT and TTS synthesis"
        ], PRIMARY_BLUE),
        ("⚙️ Backend & API Gateway", [
            "Node.js & Express: RESTful endpoints with high concurrency",
            "Mongoose ODM: Strict schema models and GeoJSON 2dsphere",
            "Supabase Auth: Secure phone OTP verification and session token",
            "CORS & Dotenv: Multi-environment security configuration"
        ], EMERALD_GREEN),
        ("🧠 AI & Cloud Data Layer", [
            "OpenRouter LLM: State-of-the-art agent reasoning and tool calls",
            "Custom RAG Vector Engine: Markdown knowledge chunk cosine search",
            "MongoDB Atlas: Cloud document storage with geo-indexing",
            "Hydration Cache: Optimistic updates and offline synchronization"
        ], PURPLE_ACCENT)
    ]

    for i, (title, items, color) in enumerate(tiers):
        left = Inches(0.8 + i * 4.0)
        add_card(slide7, left, Inches(1.8), Inches(3.7), Inches(4.8), CARD_BG, color)
        tb = slide7.shapes.add_textbox(left + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.3))
        tf = tb.text_frame
        tf.word_wrap = True
        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(14)
        pt.font.bold = True
        pt.font.color.rgb = color
        for it in items:
            p = tf.add_paragraph()
            p.text = "• " + it
            p.font.size = Pt(9.5)
            p.font.color.rgb = TEXT_WHITE

    # =========================================================================
    # SLIDE 8: Synchronized Persona User Journey
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide8, BG_DARK)
    add_header(slide8, "Priya & Vikram: The Synchronized 5-Step Storyboard", "USER JOURNEY", "🌟 5-Step Flow")

    journey_steps = [
        ("01. SOS Trigger", "Priya faces highway flat tyre at 9 PM; taps red SOS button on Skill-Link.", SOS_RED),
        ("02. GPS Geo-Query", "Lexi AI locks exact GPS coordinates and executes 5km radius query for mobile mechanics.", CYAN_ACCENT),
        ("03. 15s Auto-Lock", "Vikram (1.8km away) receives ticket alert and confirms in 6s; dispatch locks instantly.", AMBER_GOLD),
        ("04. Rapid On-Site Fix", "Vikram arrives on motorcycle in 12 mins with tools, replacing tyre safely at fixed rate.", EMERALD_GREEN),
        ("05. Payment & Rating", "Priya pays securely via UPI, receives 7-day guarantee invoice, and rates Vikram 5 stars.", PURPLE_ACCENT)
    ]

    for i, (st, sd, color) in enumerate(journey_steps):
        top = Inches(1.8 + i * 1.0)
        add_card(slide8, Inches(0.8), top, Inches(11.733), Inches(0.88), CARD_BG, color)
        
        # Step Number Badge
        num_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), top + Inches(0.12), Inches(1.8), Inches(0.6))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = color
        tf_n = num_box.text_frame
        p_n = tf_n.paragraphs[0]
        p_n.text = st
        p_n.font.size = Pt(11)
        p_n.font.bold = True
        p_n.font.color.rgb = TEXT_WHITE
        p_n.alignment = PP_ALIGN.CENTER

        # Step Text
        tb = slide8.shapes.add_textbox(Inches(3.0), top + Inches(0.12), Inches(9.3), Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = sd
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_WHITE

    # =========================================================================
    # SLIDE 9: Database Schema & Geo-Spatial Modeling
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide9, BG_DARK)
    add_header(slide9, "Data Modeling: Schema Design & Geo-Spatial Indexing", "DATA ARCHITECTURE", "🌍 2dsphere Query")

    schemas = [
        ("👤 Worker Model", [
            "name, phone (Unique)",
            "skills[] (e.g. Mechanic, Electrician)",
            "hourlyRate / calloutFee",
            "rating & totalReviews",
            "isAvailable & isVerified",
            "location: { type: 'Point', coordinates: [lng, lat] }",
            "Index: 2dsphere for geoNear queries"
        ], PRIMARY_BLUE),
        ("📅 Booking & SOS Model", [
            "bookingId: 'SOS-8921-X'",
            "customerId & workerId (ObjectId refs)",
            "serviceType & mode (HOME / SOS)",
            "urgency: HIGH / MEDIUM",
            "status: Pending / Dispatched / Completed",
            "location: { address, coordinates }",
            "priceLock & paymentStatus"
        ], SOS_RED),
        ("📚 RAG Knowledge Chunk", [
            "docTitle: 'Rate Card 2026'",
            "category: 'PRICING_POLICY'",
            "chunkText: 'Standard puncture rate ₹250...'",
            "embeddingVector: [0.012, -0.045, ...]",
            "tags[]: ['mechanic', 'highway', 'delhi']",
            "lastIndexed: ISODate()"
        ], PURPLE_ACCENT)
    ]

    for i, (title, items, color) in enumerate(schemas):
        left = Inches(0.8 + i * 4.0)
        add_card(slide9, left, Inches(1.8), Inches(3.7), Inches(4.8), CARD_BG, color)
        tb = slide9.shapes.add_textbox(left + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.3))
        tf = tb.text_frame
        tf.word_wrap = True
        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(14)
        pt.font.bold = True
        pt.font.color.rgb = color
        for it in items:
            p = tf.add_paragraph()
            p.text = "• " + it
            p.font.size = Pt(9.5)
            p.font.color.rgb = TEXT_WHITE

    # =========================================================================
    # SLIDE 10: Market Differentiation Matrix
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide10, BG_DARK)
    add_header(slide10, "Market Differentiation: Why Skill-Link Stands Out", "COMPETITIVE ANALYSIS", "🏆 Market Leader")

    matrix = [
        ("Emergency Highway Roadside SOS", "❌ Not Supported (Home-only)", "✅ 15-Min GPS Priority Dispatch"),
        ("Vernacular Voice AI (Hindi / English)", "❌ Complex text forms only", "✅ Sahayak Voice (Zero-typing)"),
        ("Autonomous AI Agent (RAG + Tools)", "❌ Simple canned FAQ bots", "✅ Lexi AI with real-time tool execution"),
        ("Interactive DIY Self-Troubleshooting", "❌ None (forces booking)", "✅ QuickFix Step-by-Step Diagnostic Trees"),
        ("Verified Direct Brand Helplines", "❌ Hidden / ad-driven results", "✅ 50+ Verified Official Toll-Free Directory"),
        ("Worker Commission Burden", "⚠️ High (20% – 35% commission)", "✅ Fair, Low-Friction (0% – 5% model)")
    ]

    add_card(slide10, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8), CARD_BG, CARD_BORDER)
    tb_m = slide10.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.3))
    tf_m = tb_m.text_frame
    tf_m.word_wrap = True

    # Header row
    p_th = tf_m.paragraphs[0]
    p_th.text = f"{'FEATURE':<38}{'LEGACY PLATFORMS':<38}{'SKILL-LINK PLATFORM'}"
    p_th.font.size = Pt(11)
    p_th.font.bold = True
    p_th.font.color.rgb = CYAN_ACCENT

    for f, l, s in matrix:
        p = tf_m.add_paragraph()
        p.text = f"{f:<38}{l:<38}{s}"
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_WHITE

    # =========================================================================
    # SLIDE 11: Real-Time Live Demo Simulation
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide11, BG_DARK)
    add_header(slide11, "Real-Time Interactive Platform Execution Simulation", "LIVE DEMO", "🎮 Simulation Mode")

    sim_cards = [
        ("🚨 Scenario 1: Highway SOS Trigger", [
            "Priya reports flat tyre at Km 32 on NH-48.",
            "High-Accuracy GPS pins Lat 28.5355, Lng 77.3910.",
            "Broadcasts to 3 nearest mechanics within 5km.",
            "Vikram accepts in 4.2s; ETA: 11 mins locked."
        ], SOS_RED),
        ("🎙️ Scenario 2: Sahayak Hindi Voice", [
            "User speaks: 'पंखा ठीक करने के लिए इलेक्ट्रीशियन चाहिए'.",
            "Tokenizer extracts intent: { skill: 'Electrician' }.",
            "System matches 4 verified local pros nearby.",
            "Voice responds back with upfront rate card confirmation."
        ], EMERALD_GREEN),
        ("🤖 Scenario 3: Lexi RAG Autonomous Query", [
            "User asks: 'Find Electrician under ₹400'.",
            "Cosine vector similarity score: 0.94 on rate doc.",
            "Executes searchWorkers(skill='Electrician', maxPrice=400).",
            "Returns interactive booking UI card inside chat."
        ], PURPLE_ACCENT)
    ]

    for i, (title, items, color) in enumerate(sim_cards):
        left = Inches(0.8 + i * 4.0)
        add_card(slide11, left, Inches(1.8), Inches(3.7), Inches(4.8), CARD_BG, color)
        tb = slide11.shapes.add_textbox(left + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.3))
        tf = tb.text_frame
        tf.word_wrap = True
        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(13)
        pt.font.bold = True
        pt.font.color.rgb = color
        for it in items:
            p = tf.add_paragraph()
            p.text = "➔ " + it
            p.font.size = Pt(9.5)
            p.font.color.rgb = TEXT_WHITE

    # =========================================================================
    # SLIDE 12: Roadmap, Vision & Finale
    # =========================================================================
    slide12 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide12, BG_DARK)
    add_header(slide12, "Future Roadmap & Project Conclusion", "THE ROAD AHEAD", "✨ Grand Finale")

    phases = [
        ("🚀 Phase 1: Live WebSocket Map", "Turn-by-turn live GPS tracking of dispatched mobile mechanics on interactive maps with sub-second WebSocket updates.", PRIMARY_BLUE),
        ("🌐 Phase 2: 8 Indian Languages", "Extending Sahayak Voice engine to Marathi, Bengali, Tamil, Telugu, and Kannada with regional dialect acoustic models.", EMERALD_GREEN),
        ("🚗 Phase 3: IoT OBD-II Telemetry", "Direct vehicle dongle telemetry diagnostics to detect battery decay and engine trouble codes before roadside breakdown occurs.", PURPLE_ACCENT)
    ]

    for i, (title, desc, color) in enumerate(phases):
        left = Inches(0.8 + i * 4.0)
        add_card(slide12, left, Inches(1.8), Inches(3.7), Inches(3.2), CARD_BG, color)
        tb = slide12.shapes.add_textbox(left + Inches(0.2), Inches(2.0), Inches(3.3), Inches(2.8))
        tf = tb.text_frame
        tf.word_wrap = True
        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(13)
        pt.font.bold = True
        pt.font.color.rgb = color
        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(10)
        pd.font.color.rgb = TEXT_MUTED

    # Bottom Conclusion Card
    add_card(slide12, Inches(0.8), Inches(5.3), Inches(11.733), Inches(1.3), CARD_BG, AMBER_GOLD)
    tb_c = slide12.shapes.add_textbox(Inches(1.0), Inches(5.45), Inches(11.3), Inches(1.0))
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True
    p_c1 = tf_c.paragraphs[0]
    p_c1.text = "✨ Skill-Link: Empowering Workers • Safeguarding Commuters • Driving AI Inclusion"
    p_c1.font.size = Pt(15)
    p_c1.font.bold = True
    p_c1.font.color.rgb = TEXT_WHITE
    p_c1.alignment = PP_ALIGN.CENTER
    
    p_c2 = tf_c.add_paragraph()
    p_c2.text = "Thank You! Open for Questions & Live Interactive Demonstration"
    p_c2.font.size = Pt(11)
    p_c2.font.color.rgb = CYAN_ACCENT
    p_c2.alignment = PP_ALIGN.CENTER

    # Save
    output_path = os.path.join(os.path.dirname(__file__), "Skill-Link_Project_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation successfully generated at: {output_path}")

if __name__ == "__main__":
    create_presentation()
